import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

# ---------- Helper Functions ----------

def export_slides_as_images(pptx_path, output_folder):
    """
    Export each slide as an image.
    Returns a list of image paths.
    """
    prs = Presentation(pptx_path)
    slide_images = []

    for i, slide in enumerate(prs.slides, start=1):
        # Create a blank image with white background
        img_width = prs.slide_width.pt
        img_height = prs.slide_height.pt
        img = Image.new('RGB', (int(img_width), int(img_height)), color='white')

        # Loop over shapes and render pictures and text
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Save picture to BytesIO and paste onto the image
                image_stream = io.BytesIO(shape.image.blob)
                picture = Image.open(image_stream)
                # Positioning needs scaling (approx)
                left = int(shape.left.pt)
                top = int(shape.top.pt)
                img.paste(picture, (left, top))
            # Note: text rendering is complex; skipping for automated images
            # For full text export, we’d need to render fonts manually

        img_path = os.path.join(output_folder, f"slide_{i}.png")
        img.save(img_path)
        slide_images.append(img_path)

    return slide_images

def load_template(template_folder, template_name):
    path = os.path.join(template_folder, template_name)
    if not os.path.exists(path):
        print(f"Template {template_name} not found in {template_folder}")
        return ""
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

def generate_canvas_page(template_html, slide_images, output_file, one_slide_per_page=True):
    if one_slide_per_page:
        base_name = os.path.splitext(output_file)[0]
        for i, slide in enumerate(slide_images, start=1):
            page_html = template_html.replace("<!--INSERT_SLIDES_HERE-->", f'<img src="{slide}" style="width:100%; margin-bottom:16px;">')
            page_file = f"{base_name}_slide{i}.html"
            with open(page_file, "w", encoding="utf-8") as f:
                f.write(page_html)
            print(f"Generated page: {page_file}")
    else:
        slides_html = ""
        for slide in slide_images:
            slides_html += f'<img src="{slide}" style="width:100%; margin-bottom:16px;">\n'
        page_html = template_html.replace("<!--INSERT_SLIDES_HERE-->", slides_html)
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(page_html)
        print(f"Generated scrollable page: {output_file}")

# ---------- Main Script ----------

def main():
    print("=== PowerPoint to Canvas Page Converter (Auto Export Images) ===\n")
    choice = input("Convert a single PPTX file or a folder? (file/folder): ").strip().lower()

    pptx_files = []

    if choice == "file":
        pptx_path = input("Enter path to the PowerPoint file (.pptx): ").strip()
        if not os.path.exists(pptx_path):
            print("File does not exist!")
            return
        pptx_files.append(pptx_path)
    elif choice == "folder":
        folder_path = input("Enter path to the folder containing PPTX files: ").strip()
        if not os.path.exists(folder_path):
            print("Folder does not exist!")
            return
        for f in os.listdir(folder_path):
            if f.lower().endswith(".pptx"):
                pptx_files.append(os.path.join(folder_path, f))
        if not pptx_files:
            print("No PPTX files found in folder.")
            return
    else:
        print("Invalid option.")
        return

    template_folder = "LMS Templates"
    template_name = "lesson_template.html"
    template_html = load_template(template_folder, template_name)
    if not template_html:
        return

    output_root = "Canvas_Output"
    os.makedirs(output_root, exist_ok=True)

    mode = input("Generate one page per slide or all slides in one page? (one/all): ").strip().lower()
    one_slide_per_page = mode == "one"

    for pptx in pptx_files:
        ppt_name = os.path.splitext(os.path.basename(pptx))[0]
        slide_folder = os.path.join(output_root, ppt_name + "_slides")
        os.makedirs(slide_folder, exist_ok=True)

        print(f"\nProcessing: {pptx}")
        slide_images = export_slides_as_images(pptx, slide_folder)

        output_file = os.path.join(output_root, f"{ppt_name}.html")
        generate_canvas_page(template_html, slide_images, output_file, one_slide_per_page)

    print("\nConversion complete! HTML files and images are in the Canvas_Output folder.")

if __name__ == "__main__":
    main()
